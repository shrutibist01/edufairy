import os
import tempfile
import mammoth
import pptx
from PyPDF2 import PdfReader
from openai import OpenAI
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.core.files.storage import default_storage

@csrf_exempt
def summarize_file(request):
    if request.method == 'POST' and request.FILES.get('file'):
        uploaded_file = request.FILES['file']
        ext = uploaded_file.name.split('.')[-1].lower()

        # Save temporarily
        tmp_path = default_storage.save(uploaded_file.name, uploaded_file)
        tmp_file = default_storage.path(tmp_path)

        try:
            text = ""
            if ext in ['ppt', 'pptx']:
                from pptx import Presentation
                prs = Presentation(tmp_file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            elif ext == 'pdf':
                reader = PdfReader(tmp_file)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
            else:
                return JsonResponse({'error': 'Unsupported file format'}, status=400)

            # Summarize with OpenAI
            client = OpenAI()
            completion = client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are a helpful assistant that summarizes lecture slides or PDFs in detail."},
                    {"role": "user", "content": f"Summarize the following content:\n\n{text}"}
                ],
                temperature=0.7
            )

            summary = completion.choices[0].message.content
            return JsonResponse({'summary': summary})

        finally:
            if os.path.exists(tmp_file):
                os.remove(tmp_file)

    return JsonResponse({'error': 'Invalid request'}, status=400)
